"""Monthly deck generator (docs/OLD_CI_System_Build_Roadmap.md, Phase 3).

Same data as the brief (surfaces/brief.py) — one run's new developments,
grouped by action, most urgent first — rendered to a branded PowerPoint deck
per `OLD - Branding Guidelines.pdf` (light mode, Sora headlines, Inter body,
IBM Plex Mono for score readouts; see surfaces/brand.py for the transcribed
constants). Pure reads (db.queries only) — generating a deck can never
mutate data.

python-pptx can only *request* Sora/Inter/IBM Plex Mono; if the machine that
opens the deck doesn't have them installed, PowerPoint substitutes a fallback
font, same as any other generated deck.

Run: python -m surfaces.deck [--run-id RUN_ID] [--out path/to/deck.pptx]
"""
import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from config import db_path
from db.init_db import connect
from db.queries import get_run, latest_run, list_current_events, list_unscored_events, sightings_for_event
from surfaces.brand import (
    ACTION_LABELS,
    ACTION_ORDER,
    ACTION_TIER_COLOR,
    COLORS,
    FONT_BODY,
    FONT_HEADLINE,
    FONT_MONO,
    PUBLIC_BRAND,
    SHORTHAND,
    TAGLINE_PRIMARY,
    hex_to_rgb,
)

DEFAULT_OUT_PATH = Path(__file__).parent / "deck" / "deck.pptx"
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def _rgb(hex_color):
    return RGBColor(*hex_to_rgb(hex_color))


def _blank_slide(prs, bg=COLORS["bg"]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(bg)
    return slide


def _add_text(slide, left, top, width, height, text, size=18, bold=False,
              color=COLORS["text_primary"], font=FONT_BODY, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font
    run.font.color.rgb = _rgb(color)
    return box


def _add_rect(slide, left, top, width, height, fill_hex, line_hex=None):
    width = max(int(width), 1)
    height = max(int(height), 1)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_hex)
    if line_hex:
        shape.line.color.rgb = _rgb(line_hex)
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _add_score_bar(slide, left, top, label, score, bucket, color):
    _add_text(slide, left, top, Inches(5), Inches(0.3), f"{label}: {score} ({bucket})",
              size=13, bold=True, font=FONT_BODY)
    track_top = top + Inches(0.35)
    track_width = Inches(6)
    _add_rect(slide, left, track_top, track_width, Inches(0.22), COLORS["surface"], line_hex=COLORS["text_secondary"])
    fill_width = int(track_width * max(0, min(score, 100)) / 100)
    _add_rect(slide, left, track_top, fill_width, Inches(0.22), color)


def _month_label(run_info):
    if not run_info or not run_info.get("started_at"):
        return ""
    return run_info["started_at"][:10]


def _build_title_slide(prs, run_info):
    slide = _blank_slide(prs)
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(0.7), Inches(1.1), Inches(1.4))
    badge.fill.solid()
    badge.fill.fore_color.rgb = _rgb(COLORS["kennel_black"])
    badge.line.color.rgb = _rgb(COLORS["lucky_copper"])
    badge.line.width = Pt(2)
    badge.shadow.inherit = False
    p = badge.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = SHORTHAND
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.name = FONT_HEADLINE
    run.font.color.rgb = _rgb(COLORS["bone"])

    _add_text(slide, Inches(2.1), Inches(0.8), Inches(9), Inches(0.5), PUBLIC_BRAND,
              size=18, bold=True, font=FONT_HEADLINE, color=COLORS["text_secondary"])
    _add_text(slide, Inches(2.1), Inches(1.3), Inches(10), Inches(1.6),
              "Monthly Competitive Intelligence Brief", size=40, bold=True, font=FONT_HEADLINE)
    _add_text(slide, Inches(2.1), Inches(2.9), Inches(8), Inches(0.5), _month_label(run_info),
              size=16, font=FONT_MONO, color=COLORS["text_secondary"])
    _add_text(slide, Inches(2.1), Inches(6.6), Inches(9), Inches(0.5), TAGLINE_PRIMARY,
              size=14, bold=True, font=FONT_BODY, color=COLORS["signal_blue"])
    return slide


def _build_section_divider(prs, index, title, subtitle, color):
    slide = _blank_slide(prs)
    _add_text(slide, Inches(0.8), Inches(2.7), Inches(2), Inches(1.5), f"{index:02d}",
              size=60, bold=True, font=FONT_HEADLINE, color=color)
    _add_text(slide, Inches(2.8), Inches(2.8), Inches(9.5), Inches(1.1), title, size=38, bold=True, font=FONT_HEADLINE)
    _add_rect(slide, Inches(2.85), Inches(3.7), Inches(3), Pt(2.5), color)
    _add_text(slide, Inches(2.8), Inches(3.85), Inches(9), Inches(0.5), subtitle,
              size=15, font=FONT_BODY, color=COLORS["text_secondary"])
    return slide


def _build_event_slide(prs, ev, sightings):
    slide = _blank_slide(prs)
    color = ACTION_TIER_COLOR.get(ev["action"], COLORS["text_secondary"])

    badge = _add_rect(slide, Inches(10.4), Inches(0.55), Inches(2.3), Inches(0.42), color)
    p = badge.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = ACTION_LABELS.get(ev["action"], ev["action"]).upper()
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.name = FONT_BODY
    run.font.color.rgb = _rgb(COLORS["bone"])

    kicker = "CONVERGENCE SIGNAL" if ev["convergence_flag"] else "COMPETITIVE SIGNAL"
    _add_text(slide, Inches(0.7), Inches(0.55), Inches(8), Inches(0.4), kicker,
              size=13, bold=True, color=COLORS["lucky_copper"], font=FONT_BODY)
    _add_text(slide, Inches(0.7), Inches(0.9), Inches(9.7), Inches(1.0), ev["title"],
              size=26, bold=True, font=FONT_HEADLINE)
    _add_text(slide, Inches(0.7), Inches(1.8), Inches(9.7), Inches(0.4),
              f"{ev['competitor_name']} · {ev['category']} · {ev['event_date']}",
              size=13, color=COLORS["text_secondary"], font=FONT_BODY)

    _add_score_bar(slide, Inches(0.7), Inches(2.45), "Industry Impact",
                    ev["industry_score"], ev["industry_bucket"], COLORS["signal_blue"])
    _add_score_bar(slide, Inches(0.7), Inches(3.25), "Relevance to OLD",
                    ev["relevance_score"], ev["relevance_bucket"], COLORS["lucky_copper"])

    if ev["requires_cco_review"]:
        _add_text(slide, Inches(7.2), Inches(2.45), Inches(4.5), Inches(0.9),
                   "Requires CCO review before confirmation", size=12, bold=True, color=COLORS["negative_red"])

    _add_text(slide, Inches(0.7), Inches(4.15), Inches(11.9), Inches(2.1),
              ev["so_what"] or "(no narration available)", size=16, font=FONT_BODY)

    if sightings:
        src_text = "Sources: " + "; ".join(f"{s['surface']} ({s['observed_at']})" for s in sightings[:4])
    else:
        src_text = "Sources: none recorded"
    _add_text(slide, Inches(0.7), Inches(6.85), Inches(11.9), Inches(0.5), src_text,
              size=10, color=COLORS["text_secondary"], font=FONT_MONO)
    return slide


def _build_pending_slide(prs, pending):
    slide = _blank_slide(prs)
    _add_text(slide, Inches(0.8), Inches(0.7), Inches(10), Inches(0.8),
              "Pending Review", size=32, bold=True, font=FONT_HEADLINE)
    _add_text(slide, Inches(0.8), Inches(1.5), Inches(10), Inches(0.5),
              "Clustered this run but not yet scored.", size=14, color=COLORS["text_secondary"])
    lines = "\n".join(f"• {ev['title']} ({ev['competitor_id']}, {ev['event_date']})" for ev in pending)
    _add_text(slide, Inches(0.8), Inches(2.2), Inches(11), Inches(4.5), lines, size=15, font=FONT_BODY)
    return slide


def build_deck(conn, run_id=None):
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    run = get_run(conn, run_id) if run_id else latest_run(conn)
    _build_title_slide(prs, run)
    if run is None:
        return prs

    all_events = [ev for ev in list_current_events(conn) if ev["created_run_id"] == run["id"]]
    pending = [ev for ev in list_unscored_events(conn) if ev["created_run_id"] == run["id"]]

    section_index = 1
    convergence_events = [ev for ev in all_events if ev["convergence_flag"]]
    if convergence_events:
        _build_section_divider(
            prs, section_index, "CONVERGENCE SIGNALS",
            f"{len(convergence_events)} event(s) this run", COLORS["lucky_copper"],
        )
        section_index += 1
        for ev in convergence_events:
            _build_event_slide(prs, ev, sightings_for_event(conn, ev["id"]))

    by_action = {}
    for ev in all_events:
        by_action.setdefault(ev["action"], []).append(ev)

    for action in ACTION_ORDER:
        events = by_action.get(action)
        if not events:
            continue
        _build_section_divider(
            prs, section_index, ACTION_LABELS[action].upper(),
            f"{len(events)} event(s) this run", ACTION_TIER_COLOR[action],
        )
        section_index += 1
        for ev in events:
            _build_event_slide(prs, ev, sightings_for_event(conn, ev["id"]))

    if pending:
        _build_pending_slide(prs, pending)

    return prs


def write_deck(conn, run_id=None, out_path=DEFAULT_OUT_PATH):
    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs = build_deck(conn, run_id=run_id)
    prs.save(str(out_path))
    return out_path


if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("--run-id", default=None)
    parser.add_argument("--out", default=str(DEFAULT_OUT_PATH))
    args = parser.parse_args()
    connection = connect(db_path())
    try:
        written_to = write_deck(connection, run_id=args.run_id, out_path=args.out)
        print(f"Deck written to {written_to}")
    finally:
        connection.close()
