import pytest
from pptx import Presentation

from db.init_db import connect, create_schema, seed_competitors
from surfaces.deck import build_deck, write_deck


@pytest.fixture()
def conn(tmp_path):
    c = connect(tmp_path / "test.sqlite3")
    create_schema(c)
    seed_competitors(c)
    c.execute("INSERT INTO run (id, type, status, started_at) VALUES ('run-1', 'monthly', 'complete', '2026-07-24T00:00:00')")

    c.execute(
        "INSERT INTO event (id, competitor_id, title, category, pillar, event_date, convergence_flag, created_run_id) "
        "VALUES ('e-priority', 'coinbase', 'Coinbase closes Deribit acquisition', 'acquisition_ma', 'crypto', "
        "'2026-07-24', 1, 'run-1')"
    )
    c.execute(
        "INSERT INTO assessment (id, event_id, rubric_version, industry_score, industry_bucket, "
        "relevance_score, relevance_bucket, headline_score, wedge_direction, action, requires_cco_review, "
        "so_what, dimension_evidence) VALUES ('a-1', 'e-priority', 'v1', 92, 'Very High', 88, 'Very High', 92, "
        "'threatens', 'PRIORITIZE', 1, 'A market-defining consolidation.', '{}')"
    )
    c.execute("UPDATE event SET current_assessment_id = 'a-1' WHERE id = 'e-priority'")
    c.execute(
        "INSERT INTO sighting (id, competitor_id, run_id, event_id, surface, source_url, observed_at, title, raw_excerpt) "
        "VALUES ('s1', 'coinbase', 'run-1', 'e-priority', 'edgar', 'https://example.test/8k', '2026-07-24', "
        "'Coinbase 8-K', 'Deribit acquisition closed')"
    )

    c.execute(
        "INSERT INTO event (id, competitor_id, title, category, pillar, event_date, created_run_id) "
        "VALUES ('e-track', 'robinhood', 'Robinhood routine filing update', 'regulatory_filing', 'options', "
        "'2026-07-10', 'run-1')"
    )
    c.execute(
        "INSERT INTO assessment (id, event_id, rubric_version, industry_score, industry_bucket, "
        "relevance_score, relevance_bucket, headline_score, wedge_direction, action, requires_cco_review, "
        "so_what, dimension_evidence) VALUES ('a-2', 'e-track', 'v1', 45, 'Moderate', 48, 'Moderate', 48, "
        "'neutral', 'TRACK', 0, 'Standard digest item.', '{}')"
    )
    c.execute("UPDATE event SET current_assessment_id = 'a-2' WHERE id = 'e-track'")

    c.execute(
        "INSERT INTO event (id, competitor_id, title, category, event_date, created_run_id) "
        "VALUES ('e-pending', 'webull', 'Webull pending signal', 'other', '2026-07-20', 'run-1')"
    )
    c.commit()
    yield c
    c.close()


def _all_text(prs):
    chunks = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                chunks.append(shape.text_frame.text)
    return "\n".join(chunks)


def test_build_deck_produces_title_dividers_and_event_slides(conn):
    prs = build_deck(conn, run_id="run-1")
    # title + convergence divider + convergence event + PRIORITIZE divider + event
    # + TRACK divider + event + pending slide = 8
    assert len(prs.slides) == 8


def test_deck_contains_event_titles_and_so_what(conn):
    prs = build_deck(conn, run_id="run-1")
    text = _all_text(prs)
    assert "Coinbase closes Deribit acquisition" in text
    assert "A market-defining consolidation." in text
    assert "Robinhood routine filing update" in text
    assert "Standard digest item." in text


def test_deck_contains_action_badges(conn):
    prs = build_deck(conn, run_id="run-1")
    text = _all_text(prs)
    assert "PRIORITIZE" in text
    assert "TRACK" in text


def test_deck_contains_convergence_divider(conn):
    prs = build_deck(conn, run_id="run-1")
    text = _all_text(prs)
    assert "CONVERGENCE SIGNALS" in text


def test_deck_contains_cco_review_flag(conn):
    prs = build_deck(conn, run_id="run-1")
    text = _all_text(prs)
    assert "Requires CCO review" in text


def test_deck_contains_pending_review_slide(conn):
    prs = build_deck(conn, run_id="run-1")
    text = _all_text(prs)
    assert "Pending Review" in text
    assert "Webull pending signal" in text


def test_deck_title_slide_has_brand_wordmark(conn):
    prs = build_deck(conn, run_id="run-1")
    text = _all_text(prs)
    assert "ONE LUCKY DOG" in text
    assert "TRUSTED. EXPERIENCED. UNAPOLOGETIC." in text


def test_deck_widescreen_dimensions(conn):
    prs = build_deck(conn, run_id="run-1")
    assert prs.slide_width == prs.slide_width  # sanity: no exception
    assert abs(prs.slide_width / 914400 - 13.333) < 0.01
    assert abs(prs.slide_height / 914400 - 7.5) < 0.01


def test_write_deck_creates_valid_pptx_file(conn, tmp_path):
    out_path = tmp_path / "out" / "deck.pptx"
    result_path = write_deck(conn, run_id="run-1", out_path=out_path)
    assert result_path == out_path
    assert out_path.exists()
    reopened = Presentation(str(out_path))
    assert len(reopened.slides) == 8


def test_build_deck_with_no_runs_returns_title_slide_only(tmp_path):
    empty_conn = connect(tmp_path / "empty.sqlite3")
    create_schema(empty_conn)
    seed_competitors(empty_conn)
    prs = build_deck(empty_conn)
    assert len(prs.slides) == 1
    empty_conn.close()
